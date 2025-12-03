#!/usr/bin/env python3
"""
Generate PowerPoint Presentation for ECG Reconstruction Project
DATA 5000 - Carleton University
Authors: Damilola Olaiya & Mithun Mani
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Create presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
CARLETON_RED = RGBColor(0xCC, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLUE = RGBColor(0x00, 0x70, 0xC0)
GREEN = RGBColor(0x00, 0xB0, 0x50)


def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    # Authors
    auth_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.333), Inches(0.5))
    tf = auth_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Damilola Olaiya & Mithun Mani"
    p.font.size = Pt(20)
    p.font.color.rgb = CARLETON_RED
    p.alignment = PP_ALIGN.CENTER
    
    # Course info
    course_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(12.333), Inches(0.5))
    tf = course_box.text_frame
    p = tf.paragraphs[0]
    p.text = "DATA 5000 - Carleton University | December 2025"
    p.font.size = Pt(16)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    return slide


def add_section_slide(prs, title):
    """Add a section divider slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Add colored background shape
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARLETON_RED
    shape.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    return slide


def add_content_slide(prs, title, bullets, notes=""):
    """Add a standard content slide with title and bullets"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    # Underline
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(12.333), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = CARLETON_RED
    line.line.fill.background()
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12.333), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        # Handle indentation levels
        if bullet.startswith("  - "):
            p.text = bullet[4:]
            p.level = 1
            p.font.size = Pt(18)
        elif bullet.startswith("    • "):
            p.text = bullet[6:]
            p.level = 2
            p.font.size = Pt(16)
        else:
            p.text = bullet.lstrip("• -")
            p.level = 0
            p.font.size = Pt(20)
        
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(8)
    
    # Speaker notes
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes
    
    return slide


def add_two_column_slide(prs, title, left_content, right_content, notes=""):
    """Add a two-column content slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    # Underline
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(12.333), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = CARLETON_RED
    line.line.fill.background()
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(5.9), Inches(5.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item.lstrip("• -")
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(6)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(6.9), Inches(1.4), Inches(5.9), Inches(5.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item.lstrip("• -")
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(6)
    
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes
    
    return slide


def add_table_slide(prs, title, headers, rows, notes=""):
    """Add a slide with a table"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    # Underline
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(12.333), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = CARLETON_RED
    line.line.fill.background()
    
    # Table
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_width = min(11, num_cols * 2.2)
    col_width = table_width / num_cols
    
    left = Inches((13.333 - table_width) / 2)
    top = Inches(1.5)
    width = Inches(table_width)
    height = Inches(0.4 * num_rows)
    
    table = slide.shapes.add_table(num_rows, num_cols, left, top, width, height).table
    
    # Set column widths
    for i in range(num_cols):
        table.columns[i].width = Inches(col_width)
    
    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = CARLETON_RED
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    
    # Data rows
    for row_idx, row in enumerate(rows):
        for col_idx, cell_text in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = DARK_GRAY
            p.alignment = PP_ALIGN.CENTER
    
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes
    
    return slide


def add_equation_slide(prs, title, equations, explanations, notes=""):
    """Add a slide with equations"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    # Underline
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(12.333), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = CARLETON_RED
    line.line.fill.background()
    
    # Equations
    y_pos = 1.5
    for eq, expl in zip(equations, explanations):
        # Equation box
        eq_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(11), Inches(0.5))
        tf = eq_box.text_frame
        p = tf.paragraphs[0]
        p.text = eq
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = BLUE
        p.alignment = PP_ALIGN.CENTER
        
        # Explanation
        expl_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos + 0.5), Inches(11), Inches(0.4))
        tf = expl_box.text_frame
        p = tf.paragraphs[0]
        p.text = expl
        p.font.size = Pt(16)
        p.font.color.rgb = LIGHT_GRAY
        p.alignment = PP_ALIGN.CENTER
        
        y_pos += 1.2
    
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes
    
    return slide


# ============================================================================
# CREATE THE PRESENTATION
# ============================================================================

# Slide 1: Title
add_title_slide(prs, 
    "12-Lead ECG Reconstruction from Reduced Lead Sets",
    "A Hybrid Physics-Informed Deep Learning Approach")

# Slide 2: Problem Motivation
add_content_slide(prs, "The Clinical Problem", [
    "• Cardiovascular disease: #1 cause of death globally (17.9M deaths/year)",
    "• 12-lead ECG: Gold standard for cardiac diagnosis",
    "  - Detects arrhythmias, MI, conduction abnormalities, hypertrophy",
    "• The accessibility gap:",
    "  - Requires 10 electrodes + precise anatomical placement",
    "  - Needs trained technicians for proper acquisition",
    "  - Difficult in ambulances, homes, remote areas",
    "  - Consumer wearables (Apple Watch, Fitbit): only 1-2 leads",
    "• Research Question: Can we bridge this gap computationally?"
], notes="The motivation is simple: 12-lead ECGs are clinically essential but practically inaccessible. Wearables capture 1-2 leads. Can we reconstruct the full 12-lead from minimal inputs?")

# Slide 3: Our Approach
add_content_slide(prs, "Our Approach: Hybrid Physics-Informed Deep Learning", [
    "• Key Insight: The 12 leads aren't independent!",
    "",
    "• Decomposition into three groups:",
    "  - 3 Input leads (I, II, V4): Measured directly",
    "  - 4 Physics leads (III, aVR, aVL, aVF): Exact equations",
    "  - 5 DL leads (V1, V2, V3, V5, V6): Neural network",
    "",
    "• Hybrid Pipeline:",
    "  - Input (3 leads) → Physics Module (4 exact) + U-Net (5 learned) → 12 leads",
    "",
    "• Advantage: Guarantee perfect reconstruction for 4 leads via physics"
], notes="Our key insight is that not all leads need to be learned. Four limb leads can be computed exactly using 19th-century physiology. We only need deep learning for the 5 chest leads.")

# Slide 4: ECG Lead System Background
add_two_column_slide(prs, "Understanding ECG Leads",
    ["What is a 'lead'?",
     "• Not the wire—a VIEW of heart activity",
     "• Voltage difference between positions",
     "• Like viewing from different camera angles",
     "",
     "12-Lead Breakdown:",
     "• Bipolar Limb: I, II, III",
     "• Augmented Limb: aVR, aVL, aVF",
     "• Precordial (Chest): V1-V6"],
    ["Key Relationships:",
     "",
     "• Limb leads are mathematically related",
     "• Given I and II → derive III, aVR, aVL, aVF",
     "",
     "• Chest leads are independent",
     "• Must be measured or reconstructed",
     "• Positioned across the chest wall",
     "• V1-V2: Right side (septum)",
     "• V4-V6: Left side (lateral wall)"],
    notes="A lead isn't a wire—it's a viewpoint. The 12-lead ECG gives us 12 angles. Critically, the 6 limb leads are mathematically related, while chest leads capture unique horizontal plane information.")

# Slide 5: Physics Module
add_equation_slide(prs, "Physics Component: Exact Lead Derivation",
    ["Lead III = Lead II − Lead I",
     "aVR = −(Lead I + Lead II) / 2",
     "aVL = Lead I − Lead II / 2",
     "aVF = Lead II − Lead I / 2"],
    ["Einthoven's Law (1903)",
     "Goldberger's Equation",
     "Goldberger's Equation", 
     "Goldberger's Equation"],
    notes="These equations come from Kirchhoff's voltage law applied to the body as a conductor. Given leads I and II, we compute the other four limb leads with perfect accuracy.")

# Slide 6: Physics Implications
add_content_slide(prs, "Physics Module: Implications", [
    "• Given I and II → derive III, aVR, aVL, aVF exactly",
    "",
    "• Properties:",
    "  - Zero learnable parameters",
    "  - Guaranteed r = 1.0 correlation",
    "  - No training required",
    "  - Works for ALL patients, ALL pathologies",
    "",
    "• This is the power of physics-informed ML:",
    "  - Eliminate learned error where deterministic solutions exist",
    "  - Use ML only where physics is insufficient",
    "",
    "• Code: src/physics.py"
], notes="For the physics-based leads, we achieve perfect reconstruction by definition. This is the key advantage of our hybrid approach—we guarantee perfect results for 4 of 12 leads without any learning.")

# Slide 7: The DL Challenge
add_content_slide(prs, "The Deep Learning Challenge: Chest Leads", [
    "• Problem: Chest leads V1-V6 have NO closed-form relationship to limb leads",
    "",
    "• Our Task:",
    "  - Input: 3 leads (I, II, V4) → tensor shape [B, 3, 5000]",
    "  - Output: 5 leads (V1, V2, V3, V5, V6) → tensor shape [B, 5, 5000]",
    "  - Formulation: Sequence-to-sequence regression",
    "",
    "• Why V4 as input?",
    "  - Central chest position",
    "  - Good correlation with V3, V5 (neighbors)",
    "",
    "• Challenge: V1 and V2 are far from V4 → hard to reconstruct"
], notes="Chest leads can't be derived mathematically—they measure unique electrical activity from different chest positions. We chose V4 as input because it's central, but this creates a challenge for V1/V2 reconstruction.")

# Slide 8: U-Net Architecture
add_content_slide(prs, "Neural Network: 1D U-Net Architecture", [
    "• Architecture: Encoder-Decoder with Skip Connections",
    "",
    "• Encoder Path (Downsampling):",
    "  - Conv1D blocks: 3 → 64 → 128 → 256 → 512 channels",
    "  - MaxPool1D(2) reduces temporal dimension: 5000 → 2500 → 1250 → 625 → 312",
    "",
    "• Bottleneck: 1024 channels, captures multi-beat context",
    "",
    "• Decoder Path (Upsampling):",
    "  - ConvTranspose1D restores resolution",
    "  - Skip connections concatenate encoder features",
    "  - Channels: 1024 → 512 → 256 → 128 → 64 → 5",
    "",
    "• Output: [B, 5, 5000] (V1, V2, V3, V5, V6)"
], notes="We use a 1D U-Net—classic encoder-decoder with skip connections, adapted for temporal signals. The encoder compresses, capturing multi-scale features. Skip connections preserve fine-grained temporal details like QRS spikes.")

# Slide 9: Architecture Decisions Table
add_table_slide(prs, "Architecture Design Decisions",
    ["Design Choice", "Decision", "Rationale"],
    [["Architecture", "1D U-Net", "Skip connections preserve QRS spikes"],
     ["Downsampling", "MaxPool1d(2)", "Preserves sharp peaks"],
     ["Upsampling", "ConvTranspose1d", "Learnable upsampling"],
     ["Normalization", "BatchNorm", "Stable training"],
     ["Kernel Size", "3", "Captures 6ms local morphology"],
     ["Dropout", "0.2", "Regularization"],
     ["Depth", "4 levels", "Full 10s receptive field"]],
    notes="U-Net is ideal because skip connections preserve sharp QRS complexes. We chose U-Net over Transformers for computational efficiency—self-attention on 5000 timesteps is expensive.")

# Slide 10: Model Comparison
add_table_slide(prs, "Ablation: Shared vs Lead-Specific Decoder",
    ["Metric", "Shared Decoder", "Lead-Specific", "Winner"],
    [["Parameters", "17.1M", "40.8M", "Shared ✓"],
     ["V1 Correlation", "0.726", "0.734", "Lead-Specific"],
     ["V2 Correlation", "0.683", "0.639", "Shared ✓"],
     ["V3 Correlation", "0.765", "0.720", "Shared ✓"],
     ["V5 Correlation", "0.824", "0.772", "Shared ✓"],
     ["V6 Correlation", "0.723", "0.671", "Shared ✓"],
     ["DL Mean", "0.744", "0.707", "Shared ✓"]],
    notes="Counter-intuitively, the simpler shared decoder wins on 4/5 leads despite 2.4x fewer parameters. With limited input information, parameter sharing acts as regularization.")

# Slide 11: Training Configuration
add_two_column_slide(prs, "Training Configuration",
    ["Loss Function:",
     "MSE on 5 DL leads only",
     "L = (1/5) Σ MSE(ŷ_k, y_k)",
     "",
     "Optimizer: AdamW",
     "• Learning Rate: 3×10⁻⁴",
     "• Weight Decay: 1×10⁻⁵",
     "• Batch Size: 64"],
    ["Training Details:",
     "• Epochs: 150",
     "• Scheduler: ReduceLROnPlateau",
     "• Mixed Precision: FP16",
     "",
     "Hardware:",
     "• GPU: NVIDIA A100 (40GB)",
     "• Training Time: 87 minutes",
     "",
     "Code: run_training.py"],
    notes="We use MSE loss on the 5 DL leads only—physics leads are exact and don't need training. AdamW with weight decay provides L2 regularization.")

# Slide 12: Dataset
add_two_column_slide(prs, "Dataset: PTB-XL",
    ["Source: PhysioNet (public)",
     "",
     "• Total ECGs: 21,837",
     "• Unique Patients: 18,885",
     "• Sampling Rate: 500 Hz",
     "• Duration: 10 seconds",
     "• Samples per Lead: 5,000"],
    ["Diagnostic Labels (SNOMED-CT):",
     "",
     "• SR: Sinus Rhythm (normal)",
     "• MI: Myocardial Infarction",
     "• AF: Atrial Fibrillation",
     "• LBBB/RBBB: Bundle Branch Block",
     "• LVH: Left Ventricular Hypertrophy",
     "",
     "Files: data/processed_full/"],
    notes="PTB-XL is the largest publicly available clinical ECG dataset. Each recording is 10 seconds at 500 Hz.")

# Slide 13: Data Splitting
add_content_slide(prs, "Patient-Wise Data Splitting (Critical!)", [
    "• The Problem:",
    "  - Some patients have multiple ECGs",
    "  - Record-wise splitting → same patient in train AND test",
    "  - Result: Data leakage → artificially inflated metrics",
    "",
    "• Our Solution:",
    "  - Split by PATIENT ID, not by record",
    "  - Each patient appears in only ONE split",
    "",
    "• Split Statistics:",
    "  - Train: 14,363 samples (70%)",
    "  - Validation: 3,086 samples (15%)",
    "  - Test: 3,086 samples (15%)",
    "",
    "• This is stricter but gives HONEST generalization estimates"
], notes="This is a critical methodological point. Many prior works use record-wise splitting, which leaks patient information. We ensure each patient appears in only ONE split.")

# Slide 14: Evaluation Metrics
add_table_slide(prs, "Evaluation Metrics",
    ["Metric", "Formula", "Interpretation", "Target"],
    [["MAE", "|y - ŷ| / N", "Amplitude error (lower better)", "< 0.05 mV"],
     ["Pearson r", "Cov(Y,Ŷ) / (σ_Y × σ_Ŷ)", "Shape similarity [-1,1]", "> 0.90"],
     ["SNR (dB)", "10·log₁₀(||y||² / ||y-ŷ||²)", "Signal quality (higher better)", "> 20 dB"]],
    notes="We use three complementary metrics. MAE captures amplitude error, correlation captures waveform shape similarity, and SNR gives a global signal quality measure.")

# Slide 15: Results - Physics Leads
add_table_slide(prs, "Results: Physics-Based Leads (Exact)",
    ["Lead", "MAE", "Correlation", "SNR"],
    [["III", "0.000", "1.000", "∞"],
     ["aVR", "0.000", "1.000", "∞"],
     ["aVL", "0.000", "1.000", "∞"],
     ["aVF", "0.000", "1.000", "∞"]],
    notes="For the physics-based leads, we achieve perfect reconstruction by definition. This is the key advantage of our hybrid approach.")

# Slide 16: Results - DL Leads
add_table_slide(prs, "Results: Deep Learning Leads",
    ["Lead", "MAE", "Correlation", "SNR (dB)"],
    [["V1", "0.036", "0.726", "17.9"],
     ["V2", "0.041", "0.683 (worst)", "17.1"],
     ["V3", "0.036", "0.765", "17.8"],
     ["V5", "0.032", "0.824 (best)", "18.7"],
     ["V6", "0.038", "0.723", "17.2"],
     ["Mean", "0.037", "0.744", "17.8"]],
    notes="V5 performs best (closest to V4 input). V2 performs worst (furthest from V4). All leads below clinical target of r > 0.90.")

# Slide 17: Overall Performance
add_content_slide(prs, "Aggregate 12-Lead Performance", [
    "• Overall Metrics (all 12 leads):",
    "  - Correlation: 0.893",
    "  - MAE: 0.0153",
    "  - SNR: 62.3 dB",
    "",
    "• But this is MISLEADING!",
    "  - Inflated by 3 input leads (trivially r=1.0)",
    "  - Inflated by 4 physics leads (exactly r=1.0)",
    "",
    "• Honest metric: DL leads average = 0.744",
    "",
    "• This is why we report per-component metrics!",
    "",
    "• File: models/model_comparison_metrics.json"
], notes="The overall 0.893 looks good but is misleading. It's inflated by input and physics leads. The honest measure of our learning is the DL leads average: 0.744.")

# Slide 18: Information Bottleneck
add_table_slide(prs, "Information Bottleneck Analysis",
    ["Target Lead", "Best Input Corr.", "Our Model r", "Analysis"],
    [["V1", "0.49 (Lead I)", "0.726", "Model learns +0.24"],
     ["V2", "0.36 (V4)", "0.683", "Model learns +0.32"],
     ["V3", "0.71 (V4)", "0.765", "Near ceiling"],
     ["V5", "0.79 (V4)", "0.824", "Near ceiling"],
     ["V6", "0.69 (Lead I)", "0.723", "Model learns +0.03"]],
    notes="V2's ground truth correlation with V4 is only 0.36. No model can perfectly reconstruct V2 from V4 because the information isn't there. Our 0.683 is actually impressive given this constraint.")

# Slide 19: SOTA Comparison
add_table_slide(prs, "Comparison with State-of-the-Art",
    ["Method", "Year", "Chest r", "Data Split"],
    [["Linear (Frank)", "1970s", "0.70-0.75", "N/A"],
     ["CNN (Sohn)", "2020", "0.85", "Record-wise"],
     ["LSTM (Lee)", "2021", "0.88", "Record-wise"],
     ["Transformer", "2023", "0.90", "Record-wise"],
     ["Ours", "2025", "0.744", "Patient-wise"]],
    notes="We're below SOTA on chest leads. But we use patient-wise splitting while most prior work uses record-wise—their metrics are likely inflated by data leakage.")

# Slide 20: Why the Gap?
add_content_slide(prs, "Gap Analysis: Why Below SOTA?", [
    "• Three factors explain the performance gap:",
    "",
    "1. Patient-wise splitting (stricter, more honest)",
    "   - Prior work often uses record-wise splits",
    "   - Data leakage inflates their metrics",
    "",
    "2. Input lead choice (I, II, V4)",
    "   - V4 has low correlation with V1/V2 (0.36-0.49)",
    "   - Information bottleneck limits reconstruction",
    "",
    "3. No data augmentation or pretrained weights",
    "",
    "• Our advantage: Physics guarantees + reproducible evaluation"
], notes="The gap is primarily explained by our stricter evaluation methodology and the information bottleneck from our input lead choice.")

# Slide 21: Key Findings
add_content_slide(prs, "Key Findings", [
    "1. Physics constraints work perfectly",
    "   - Limb leads: r = 1.0 guaranteed, eliminates error for 4/12 leads",
    "",
    "2. Simpler models win with limited input",
    "   - Shared decoder (17.1M) > Lead-specific (40.8M)",
    "   - Parameter sharing = implicit regularization",
    "",
    "3. Information bottleneck is the fundamental limit",
    "   - V4 ↔ V2 correlation only 0.36 in ground truth",
    "   - Model can't reconstruct what input doesn't contain",
    "",
    "4. Input selection > architecture",
    "   - Adding V2 as input would directly improve V2 reconstruction"
], notes="Five key takeaways: physics works, simpler models win, information bottleneck limits us, input selection matters most.")

# Slide 22: Limitations
add_content_slide(prs, "Limitations", [
    "1. Single Dataset",
    "   - Only PTB-XL tested",
    "   - Unknown generalization to other populations, devices, hospitals",
    "",
    "2. Resting ECGs Only",
    "   - No stress/exercise ECGs tested",
    "",
    "3. Input Lead Dependency",
    "   - Results specific to (I, II, V4) configuration",
    "",
    "4. No Downstream Validation",
    "   - Signal metrics only, no diagnostic classification test",
    "",
    "5. No Uncertainty Quantification",
    "   - Point estimates only, no confidence intervals"
], notes="Important limitations: single dataset, no downstream diagnostic validation, no uncertainty quantification.")

# Slide 23: Future Work
add_content_slide(prs, "Future Directions", [
    "1. Input Lead Optimization",
    "   - Test: (I, II, V2), (I, II, V3), (I, II, V2+V4)",
    "",
    "2. Uncertainty Quantification",
    "   - Add probabilistic head (VAE, MC Dropout)",
    "",
    "3. External Validation",
    "   - Chapman-Shaoxing (China), MIMIC-IV-ECG (US)",
    "",
    "4. Downstream Validation",
    "   - Multi-label classification on reconstructed ECGs",
    "   - Measure ΔAUROC for MI, AF, LVH detection",
    "",
    "5. Foundation Model Approach",
    "   - Self-supervised pretraining on millions of ECGs"
], notes="Most important next step is testing different input leads. We also need downstream validation to prove clinical utility.")

# Slide 24: Repository
add_content_slide(prs, "Code Repository", [
    "• GitHub: github.com/whiteblaze143/DATA_5000",
    "",
    "• Key Files:",
    "  - run_training.py: Main entry point",
    "  - src/train.py: Training loop",
    "  - src/physics.py: Einthoven/Goldberger equations",
    "  - src/models/unet_1d.py: U-Net architecture",
    "  - src/evaluation.py: Metrics computation",
    "",
    "• Data: data/processed_full/",
    "",
    "• Models: models/overnight_full/best_model.pt",
    "",
    "• Results: models/model_comparison_metrics.json"
], notes="All code is available on GitHub. The key files are in src/ for the model and training, data/ for preprocessing, and models/ for trained checkpoints.")

# Slide 25: Conclusion
add_content_slide(prs, "Conclusion", [
    "• What we built:",
    "  - Hybrid physics-informed deep learning system",
    "  - Reconstructs 12-lead ECG from 3 measured leads",
    "",
    "• What we achieved:",
    "  - ✓ Perfect limb lead reconstruction (r = 1.0) via physics",
    "  - △ Moderate chest lead reconstruction (r = 0.744) via U-Net",
    "  - ✓ Overall 12-lead correlation: 0.893",
    "",
    "• Key insight:",
    "  - Input lead selection is MORE important than model architecture",
    "",
    "• Clinical potential:",
    "  - Enable wearables to provide near-12-lead capability"
], notes="Our hybrid approach successfully combines physics and deep learning. The physics component works perfectly. The DL component is limited by the information bottleneck.")

# Slide 26: Questions
add_section_slide(prs, "Questions?")

# Final slide with contact
add_content_slide(prs, "Thank You!", [
    "",
    "• Repository: github.com/whiteblaze143/DATA_5000",
    "",
    "• Contact:",
    "  - Damilola Olaiya: damilolaolaiya@cmail.carleton.ca",
    "  - Mithun Mani: mithunmani@cmail.carleton.ca",
    "",
    "• Acknowledgments:",
    "  - DATA 5000 Course Instructors & TAs",
    "  - PhysioNet for PTB-XL dataset access"
])

# Save the presentation
output_path = os.path.join(os.path.dirname(os.path.dirname(__file__)), "ECG_Reconstruction_Presentation.pptx")
prs.save(output_path)
print(f"✓ Presentation saved to: {output_path}")
print(f"  Total slides: {len(prs.slides)}")
