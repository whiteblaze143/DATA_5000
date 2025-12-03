#!/usr/bin/env python3
"""
Final PowerPoint Presentation for ECG Reconstruction Project
DATA 5000 - Carleton University
Authors: Damilola Olaiya & Mithun Mani

CRITICALLY REFRAMED:
- Computational efficiency as primary contribution
- Rigorous experimental methodology (frozen HPs, statistical tests)
- Honest assessment of results vs SOTA
- Focus on methodological contributions
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Presentation setup
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
PRIMARY = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0xE9, 0x4D, 0x60)
ACCENT2 = RGBColor(0x0F, 0x4C, 0x75)
ACCENT3 = RGBColor(0x3E, 0xBA, 0xC2)
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_LIGHT = RGBColor(0xE0, 0xE0, 0xE0)
SUCCESS = RGBColor(0x4E, 0xC9, 0xB0)
WARNING = RGBColor(0xFF, 0xCC, 0x00)

FIGURES = "/home/mithunmanivannan/DATA_5000/DATA_5000/docs/figures"
MODELS = "/home/mithunmanivannan/DATA_5000/DATA_5000/models"


def bg(slide, color=PRIMARY):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    spTree = slide.shapes._spTree
    sp = s._element
    spTree.remove(sp)
    spTree.insert(2, sp)

def bar(slide, y=0, h=0.06, color=ACCENT):
    b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(y), Inches(13.333), Inches(h))
    b.fill.solid()
    b.fill.fore_color.rgb = color
    b.line.fill.background()

def title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    bar(slide, 0, 0.12, ACCENT)
    
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(2.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.alignment = PP_ALIGN.CENTER
    
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(4.1), Inches(4.3), Inches(0.02)).fill.solid()
    
    sb = slide.shapes.add_textbox(Inches(0.8), Inches(4.3), Inches(11.7), Inches(0.8))
    p = sb.text_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(20)
    p.font.color.rgb = ACCENT3
    p.alignment = PP_ALIGN.CENTER
    
    ab = slide.shapes.add_textbox(Inches(0.8), Inches(5.4), Inches(11.7), Inches(0.5))
    p = ab.text_frame.paragraphs[0]
    p.text = "Damilola Olaiya  •  Mithun Mani"
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_LIGHT
    p.alignment = PP_ALIGN.CENTER
    
    cb = slide.shapes.add_textbox(Inches(0.8), Inches(6.1), Inches(11.7), Inches(0.4))
    p = cb.text_frame.paragraphs[0]
    p.text = "DATA 5000  |  Carleton University  |  December 2025"
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_LIGHT
    p.alignment = PP_ALIGN.CENTER
    
    bar(slide, 7.38, 0.12, ACCENT)

def section_slide(prs, num, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, ACCENT)
    
    nb = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(11.7), Inches(1.5))
    p = nb.text_frame.paragraphs[0]
    p.text = f"0{num}" if num < 10 else str(num)
    p.font.size = Pt(120)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.alignment = PP_ALIGN.CENTER
    
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(11.7), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = title.upper()
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.alignment = PP_ALIGN.CENTER

def content_slide(prs, title, bullets, fig=None, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    bar(slide, 0, 0.06, ACCENT)
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.65))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    
    cw = 6 if fig and os.path.exists(fig) else 12
    cb = slide.shapes.add_textbox(Inches(0.5), Inches(0.95), Inches(cw), Inches(6))
    tf = cb.text_frame
    tf.word_wrap = True
    
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        txt = b.strip().lstrip("•-▸ ")
        indent = 1 if b.startswith("  ") else 0
        
        if txt.startswith("✓"): p.font.color.rgb = SUCCESS
        elif txt.startswith("△") or txt.startswith("⚠"): p.font.color.rgb = WARNING
        elif txt.startswith("✗"): p.font.color.rgb = ACCENT
        else: p.font.color.rgb = TEXT_LIGHT
        
        p.text = ("▸ " + txt) if indent == 0 and txt else ("   " + txt)
        p.level = indent
        p.font.size = Pt(16 if indent == 0 else 14)
        p.space_after = Pt(5)
    
    if fig and os.path.exists(fig):
        slide.shapes.add_picture(fig, Inches(6.7), Inches(1.1), width=Inches(6.2))
    
    if notes:
        slide.notes_slide.notes_text_frame.text = notes

def figure_slide(prs, title, fig, caption=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    bar(slide, 0, 0.06, ACCENT)
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    
    if os.path.exists(fig):
        slide.shapes.add_picture(fig, Inches(0.7), Inches(0.9), width=Inches(11.9))
    
    if caption:
        cb = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12), Inches(0.5))
        p = cb.text_frame.paragraphs[0]
        p.text = caption
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_LIGHT
        p.alignment = PP_ALIGN.CENTER

def table_slide(prs, title, headers, rows, highlight=-1):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    bar(slide, 0, 0.06, ACCENT)
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    
    nr, nc = len(rows) + 1, len(headers)
    tw = min(12, nc * 2.1)
    cw = tw / nc
    
    left = Inches((13.333 - tw) / 2)
    table = slide.shapes.add_table(nr, nc, left, Inches(1.1), Inches(tw), Inches(0.4 * nr)).table
    
    for i in range(nc):
        table.columns[i].width = Inches(cw)
    
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_WHITE
        p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            cell.fill.solid()
            if ri == highlight:
                cell.fill.fore_color.rgb = RGBColor(0x1A, 0x4A, 0x3A)
            else:
                cell.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x3A) if ri % 2 == 0 else RGBColor(0x1C, 0x1C, 0x32)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.font.color.rgb = TEXT_LIGHT
            p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

def metrics_slide(prs, title, metrics):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    bar(slide, 0, 0.06, ACCENT)
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    
    n = len(metrics)
    cw = 3.0
    sp = (13.333 - n * cw) / (n + 1)
    
    for i, (label, val, color) in enumerate(metrics):
        x = sp + i * (cw + sp)
        
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.6), Inches(cw), Inches(4))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0x20, 0x20, 0x38)
        card.line.color.rgb = color
        card.line.width = Pt(2)
        
        vb = slide.shapes.add_textbox(Inches(x), Inches(2.4), Inches(cw), Inches(1.2))
        p = vb.text_frame.paragraphs[0]
        p.text = val
        p.font.size = Pt(42)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER
        
        lb = slide.shapes.add_textbox(Inches(x), Inches(3.8), Inches(cw), Inches(1.4))
        tf = lb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_LIGHT
        p.alignment = PP_ALIGN.CENTER

def comparison_slide(prs, title, l_label, l_val, r_label, r_val, winner="left"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    bar(slide, 0, 0.06, ACCENT)
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    
    lc = SUCCESS if winner == "left" else RGBColor(0x50, 0x50, 0x70)
    lcard = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.2), Inches(5.5), Inches(5))
    lcard.fill.solid()
    lcard.fill.fore_color.rgb = RGBColor(0x20, 0x20, 0x38)
    lcard.line.color.rgb = lc
    lcard.line.width = Pt(3)
    
    llb = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(5.5), Inches(0.5))
    p = llb.text_frame.paragraphs[0]
    p.text = l_label
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = lc
    p.alignment = PP_ALIGN.CENTER
    
    lvb = slide.shapes.add_textbox(Inches(0.7), Inches(2.6), Inches(5.5), Inches(2.5))
    tf = lvb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = l_val
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.alignment = PP_ALIGN.CENTER
    
    vsb = slide.shapes.add_textbox(Inches(6), Inches(3.2), Inches(1.3), Inches(1))
    p = vsb.text_frame.paragraphs[0]
    p.text = "vs"
    p.font.size = Pt(26)
    p.font.color.rgb = TEXT_LIGHT
    p.alignment = PP_ALIGN.CENTER
    
    rc = SUCCESS if winner == "right" else RGBColor(0x50, 0x50, 0x70)
    rcard = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.1), Inches(1.2), Inches(5.5), Inches(5))
    rcard.fill.solid()
    rcard.fill.fore_color.rgb = RGBColor(0x20, 0x20, 0x38)
    rcard.line.color.rgb = rc
    rcard.line.width = Pt(3)
    
    rlb = slide.shapes.add_textbox(Inches(7.1), Inches(1.5), Inches(5.5), Inches(0.5))
    p = rlb.text_frame.paragraphs[0]
    p.text = r_label
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = rc
    p.alignment = PP_ALIGN.CENTER
    
    rvb = slide.shapes.add_textbox(Inches(7.1), Inches(2.6), Inches(5.5), Inches(2.5))
    tf = rvb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = r_val
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.alignment = PP_ALIGN.CENTER
    
    if winner:
        bx = Inches(4.7) if winner == "left" else Inches(11.1)
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, bx, Inches(0.9), Inches(1), Inches(1))
        badge.fill.solid()
        badge.fill.fore_color.rgb = SUCCESS
        badge.line.fill.background()
        cb = slide.shapes.add_textbox(bx, Inches(1), Inches(1), Inches(0.9))
        p = cb.text_frame.paragraphs[0]
        p.text = "✓"
        p.font.size = Pt(30)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        p.alignment = PP_ALIGN.CENTER

def equation_slide(prs, title, equations):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    bar(slide, 0, 0.06, ACCENT)
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    
    y = 1.1
    for name, formula in equations:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(y), Inches(11.5), Inches(0.9))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0x20, 0x20, 0x38)
        box.line.color.rgb = ACCENT2
        box.line.width = Pt(1)
        
        nb = slide.shapes.add_textbox(Inches(1.2), Inches(y + 0.1), Inches(4), Inches(0.3))
        p = nb.text_frame.paragraphs[0]
        p.text = name
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = ACCENT3
        
        fb = slide.shapes.add_textbox(Inches(1.2), Inches(y + 0.42), Inches(10), Inches(0.4))
        p = fb.text_frame.paragraphs[0]
        p.text = formula
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        
        y += 1.05


# ============================================================================
# BUILD PRESENTATION
# ============================================================================

# SLIDE 1: Title
title_slide(prs, 
    "Efficient 12-Lead ECG Reconstruction\nfrom Reduced Lead Sets",
    "A Computationally Lightweight Physics-Informed Approach")

# SLIDE 2: Section - Motivation
section_slide(prs, 1, "Motivation")

# SLIDE 3: Clinical Problem
content_slide(prs, "The Clinical Need for Accessible ECG", [
    "Cardiovascular Disease: #1 cause of mortality globally",
    "  17.9 million deaths annually",
    "",
    "12-Lead ECG: Gold standard for cardiac diagnosis",
    "  Detects MI, arrhythmias, conduction abnormalities",
    "",
    "The Accessibility Gap:",
    "  10 electrodes + precise anatomical placement",
    "  Trained technicians required",
    "  Not feasible: ambulances, homes, wearables",
    "",
    "Consumer Wearables: 1-3 leads only",
    "  Apple Watch, Kardia, Fitbit"
], fig=f"{FIGURES}/problem_visualization.png")

# SLIDE 4: SOTA - The Computational Problem
content_slide(prs, "Current SOTA: Computationally Expensive", [
    "Transformer-Based Methods (2023):",
    "  Self-attention: O(n²) for 5000 samples",
    "  100M+ parameters typical",
    "  Multi-GPU training required",
    "",
    "LSTM/BiLSTM Methods (2021):",
    "  60-80M parameters",
    "  Sequential processing = slow inference",
    "",
    "Key Problem:",
    "  These methods reconstruct ALL 9 missing leads",
    "  But 4 leads don't need learning!",
    "",
    "Our Insight: Exploit physiological constraints"
])

# SLIDE 5: Our Key Insight
content_slide(prs, "Key Insight: Physics Reduces the Problem by 44%", [
    "The 12 leads are NOT independent!",
    "",
    "4 Leads are Deterministic (Physics):",
    "  III = II - I        (Einthoven 1903)",
    "  aVR = -(I + II) / 2",
    "  aVL = I - II / 2",
    "  aVF = II - I / 2",
    "",
    "Only 5 Leads Need Deep Learning:",
    "  V1, V2, V3, V5, V6",
    "",
    "✓ Zero parameters for 4 leads",
    "✓ Perfect reconstruction guaranteed",
    "✓ Reduced model capacity needed"
], fig=f"{FIGURES}/approach_diagram.png")

# SLIDE 6: Section - Methods
section_slide(prs, 2, "Methodology")

# SLIDE 7: Hybrid Architecture
content_slide(prs, "Hybrid Architecture: Physics + Lightweight DL", [
    "Input: 3 measured leads (I, II, V4)",
    "",
    "Physics Module:",
    "  Einthoven's & Goldberger's laws",
    "  4 leads, zero cost, r = 1.0",
    "",
    "Deep Learning Module:",
    "  1D U-Net (encoder-decoder with skip connections)",
    "  17.1M parameters",
    "  5 chest leads: V1, V2, V3, V5, V6",
    "",
    "Output: Complete 12-lead ECG",
    "",
    "Total: 6x fewer params than SOTA transformers"
], fig=f"{FIGURES}/architecture_diagram_clean.png")

# SLIDE 8: Physics Equations
equation_slide(prs, "Physics Module: Zero-Cost Perfect Reconstruction", [
    ("Einthoven's Law (1903)", "Lead III = Lead II − Lead I"),
    ("Goldberger's Equations", "aVR = −(I + II) / 2"),
    ("", "aVL = I − II / 2"),
    ("", "aVF = II − I / 2"),
])

# SLIDE 9: Physics Results
metrics_slide(prs, "Physics Module: Guaranteed r = 1.0", [
    ("Lead III", "1.000", SUCCESS),
    ("aVR, aVL, aVF", "1.000", SUCCESS),
    ("Compute Cost", "~0", ACCENT3),
])

# SLIDE 10: Experimental Methodology
content_slide(prs, "Rigorous Experimental Methodology", [
    "Frozen Hyperparameters (validated via LR sweep):",
    "  lr = 3×10⁻⁴, batch = 64, epochs = 150",
    "  AdamW (weight_decay = 1×10⁻⁴)",
    "  seed = 42",
    "",
    "Model Variants Tested:",
    "  1. Baseline: Shared encoder-decoder (17.1M)",
    "  2. Hybrid: Shared + per-lead heads (17.13M)",
    "  3. Lead-Specific: Per-lead decoders (40.8M)",
    "",
    "Statistical Framework:",
    "  Paired t-test, Wilcoxon signed-rank",
    "  Cohen's d effect size",
    "  Bootstrap 95% CI",
    "  Bonferroni correction"
])

# SLIDE 11: Model Ablation
comparison_slide(prs, 
    "Ablation: Shared vs Lead-Specific Decoder",
    "Shared Decoder (Ours)", "17.1M params\nr = 0.744\n(5.3% better)",
    "Lead-Specific", "40.8M params\nr = 0.707",
    winner="left")

# SLIDE 12: Why Shared Wins
content_slide(prs, "Why Simpler Models Win (Counter-Intuitive)", [
    "Lead-Specific (40.8M) LOST to Shared (17.1M)",
    "",
    "Statistical Evidence:",
    "  Cohen's d = 0.92 (large effect)",
    "  Bootstrap 95% CI: [0.006, 0.072]",
    "  CI excludes zero → significant",
    "",
    "Per-Lead Breakdown:",
    "  Shared wins: V1, V2, V3, V5 (4/5)",
    "  Lead-Specific wins: V6 only",
    "",
    "Interpretation:",
    "  With limited input information (3 leads),",
    "  parameter sharing provides regularization",
    "  More params → overfitting, not capacity"
])

# SLIDE 13: Section - Data
section_slide(prs, 3, "Dataset & Evaluation")

# SLIDE 14: PTB-XL Dataset
content_slide(prs, "Dataset: PTB-XL", [
    "PhysioNet Public Dataset:",
    "  21,837 12-lead ECG recordings",
    "  18,885 unique patients",
    "  500 Hz, 10 seconds (5000 samples)",
    "",
    "Diagnostic Labels (SNOMED-CT):",
    "  MI, AF, LBBB, RBBB, LVH, etc.",
    "",
    "Preprocessing:",
    "  Z-score normalization per lead",
    "  Outlier removal (2.5-97.5 percentile)"
], fig=f"{FIGURES}/dataset_characteristics.png")

# SLIDE 15: Patient-wise Splitting
content_slide(prs, "Critical: Patient-Wise Data Splitting", [
    "The Data Leakage Problem:",
    "  Some patients have multiple ECGs",
    "  Record-wise split → train/test contamination",
    "  → Inflated, unrealistic metrics",
    "",
    "Our Solution (Stricter than SOTA):",
    "  Split by PATIENT ID, not record",
    "  Each patient in exactly ONE split",
    "",
    "Split Statistics:",
    "  Train: ~15,286 (70%)",
    "  Validation: ~3,276 (15%)",
    "  Test: ~3,275 (15%)",
    "",
    "△ Most prior work uses record-wise splits!"
])

# SLIDE 16: Section - Results
section_slide(prs, 4, "Results")

# SLIDE 17: Training Curves
figure_slide(prs, "Training Convergence",
    f"{MODELS}/overnight_full/training_curves.png",
    "Smooth convergence over 150 epochs with ReduceLROnPlateau scheduler")

# SLIDE 18: DL Results Table
table_slide(prs, "Deep Learning Lead Reconstruction Results",
    ["Lead", "MAE", "Pearson r", "SNR (dB)", "Category"],
    [["V1", "0.036", "0.726", "17.9", "Right Precordial"],
     ["V2", "0.041", "0.683", "17.1", "Right Precordial"],
     ["V3", "0.036", "0.765", "17.8", "Transition"],
     ["V5", "0.032", "0.824", "18.7", "Left Precordial"],
     ["V6", "0.038", "0.723", "17.2", "Left Precordial"],
     ["Mean", "0.037", "0.744", "17.8", "—"]],
    highlight=4)

# SLIDE 19: Overall Performance
metrics_slide(prs, "Overall 12-Lead Performance", [
    ("Physics Leads\n(III, aVR, aVL, aVF)", "r = 1.0", SUCCESS),
    ("DL Leads\n(V1-V3, V5-V6)", "r = 0.744", WARNING),
    ("Overall\n12-Lead", "r = 0.893", ACCENT3),
])

# SLIDE 20: Per-Lead Analysis
figure_slide(prs, "Per-Lead Performance Analysis",
    f"{FIGURES}/per_lead_performance.png",
    "V5 best (close to V4), V2 hardest (anatomically distant)")

# SLIDE 21: Section - Analysis
section_slide(prs, 5, "Analysis")

# SLIDE 22: Efficiency Comparison
table_slide(prs, "Computational Efficiency vs SOTA",
    ["Method", "Parameters", "Complexity", "Physics", "Train Time"],
    [["Transformer", "100M+", "O(n²)", "No", "Days"],
     ["BiLSTM", "60M", "O(n)", "No", "Hours"],
     ["CNN-only", "30M", "O(n)", "No", "Hours"],
     ["Ours", "17.1M", "O(n)", "Yes (4 leads)", "87 min"]],
    highlight=3)

# SLIDE 23: SOTA Comparison
table_slide(prs, "Performance Comparison with State-of-the-Art",
    ["Method", "Year", "Chest r", "Params", "Data Split"],
    [["Linear (Frank)", "1970s", "0.70-0.75", "~0", "N/A"],
     ["CNN (Sohn)", "2020", "~0.85", "~30M", "Record"],
     ["LSTM (Lee)", "2021", "~0.88", "~60M", "Record"],
     ["Transformer", "2023", "~0.90", "100M+", "Record"],
     ["Ours", "2025", "0.744", "17.1M", "Patient"]],
    highlight=4)

# SLIDE 24: Honest Gap Analysis
content_slide(prs, "Honest Assessment: Why Below SOTA?", [
    "Our r = 0.744 vs SOTA r = 0.85-0.90",
    "",
    "Three Contributing Factors:",
    "",
    "1. Stricter Evaluation (Patient-wise)",
    "   Prior work may have data leakage",
    "",
    "2. Information Bottleneck",
    "   V4 ↔ V2 ground truth correlation = 0.36",
    "   Fundamental limit, not model failure",
    "",
    "3. No Augmentation or Pretraining",
    "   Clean baseline for future work",
    "",
    "✓ Our advantage: 6x efficiency + physics guarantees"
], fig=f"{FIGURES}/information_bottleneck_analysis.png")

# SLIDE 25: Information Bottleneck
figure_slide(prs, "Information Bottleneck: Ground Truth Correlations",
    f"{FIGURES}/lead_correlation_heatmap.png",
    "V4↔V2 correlation only 0.36 — this limits reconstruction performance fundamentally")

# SLIDE 26: Key Findings
content_slide(prs, "Key Findings", [
    "✓ Physics constraints eliminate 4 leads from learning",
    "   Zero cost, perfect accuracy, guaranteed",
    "",
    "✓ Lightweight U-Net is sufficient",
    "   17.1M params vs 100M+ transformers",
    "   87 min training on A100",
    "",
    "✓ Shared decoder beats lead-specific",
    "   Cohen's d = 0.92 (large effect)",
    "   Parameter sharing = regularization",
    "",
    "⚠ Input lead selection is crucial",
    "   V4 limits V2 reconstruction",
    "   Future: optimize input configuration"
])

# SLIDE 27: Section - Conclusion
section_slide(prs, 6, "Conclusion")

# SLIDE 28: Contributions
content_slide(prs, "Contributions", [
    "1. Computationally Efficient Architecture",
    "   6x fewer parameters than SOTA",
    "   Physics module = zero inference cost",
    "",
    "2. Rigorous Evaluation Framework",
    "   Patient-wise splitting (no leakage)",
    "   Statistical tests (Cohen's d, bootstrap CI)",
    "   Frozen hyperparameters",
    "",
    "3. Information Bottleneck Analysis",
    "   Identified fundamental limits",
    "   Guides future input selection",
    "",
    "4. Open Source Implementation",
    "   github.com/whiteblaze143/DATA_5000"
])

# SLIDE 29: Limitations
content_slide(prs, "Limitations", [
    "Single Dataset",
    "  PTB-XL only — external validation needed",
    "",
    "Input Lead Configuration",
    "  (I, II, V4) may not be optimal",
    "  V4 has low correlation with V1/V2",
    "",
    "No Downstream Validation",
    "  Signal metrics only",
    "  Classification accuracy untested",
    "",
    "No Uncertainty Quantification",
    "  Point estimates only",
    "  Clinical deployment needs confidence"
])

# SLIDE 30: Future Work
content_slide(prs, "Future Directions", [
    "Input Lead Optimization",
    "  Test (I, II, V2), (I, II, V2+V4)",
    "  Ablation for optimal configuration",
    "",
    "External Validation",
    "  Chapman-Shaoxing, MIMIC-IV-ECG",
    "",
    "Downstream Task Evaluation",
    "  Multi-label classification on reconstructed ECGs",
    "",
    "Uncertainty Quantification",
    "  Probabilistic outputs (VAE, MC Dropout)",
    "",
    "Foundation Model Exploration",
    "  Leverage pretrained ECG representations"
])

# SLIDE 31: Summary
metrics_slide(prs, "Summary", [
    ("Parameters", "17.1M\n(6x smaller)", SUCCESS),
    ("Physics Leads", "r = 1.0\n(zero cost)", SUCCESS),
    ("DL Leads", "r = 0.744\n(patient-split)", ACCENT3),
])

# SLIDE 32: Code & Resources
content_slide(prs, "Code & Resources", [
    "GitHub Repository:",
    "  github.com/whiteblaze143/DATA_5000",
    "",
    "Key Files:",
    "  run_training.py — Main entry point",
    "  src/physics.py — Einthoven/Goldberger",
    "  src/models/unet_1d.py — U-Net variants",
    "  scripts/compare_variants.py — Statistical analysis",
    "",
    "Trained Models:",
    "  models/overnight_full/best_model.pt",
    "",
    "Documentation:",
    "  docs/PROJECT_REPORT.tex"
])

# SLIDE 33: Questions
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)

qb = slide.shapes.add_textbox(Inches(0), Inches(1.3), Inches(13.333), Inches(3))
p = qb.text_frame.paragraphs[0]
p.text = "?"
p.font.size = Pt(200)
p.font.bold = True
p.font.color.rgb = ACCENT
p.alignment = PP_ALIGN.CENTER

tb = slide.shapes.add_textbox(Inches(0), Inches(4.5), Inches(13.333), Inches(1))
p = tb.text_frame.paragraphs[0]
p.text = "Questions?"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = TEXT_WHITE
p.alignment = PP_ALIGN.CENTER

# SLIDE 34: Thank You
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
bar(slide, 0, 0.1, ACCENT)
bar(slide, 7.4, 0.1, ACCENT)

tyb = slide.shapes.add_textbox(Inches(0), Inches(1.8), Inches(13.333), Inches(1.2))
p = tyb.text_frame.paragraphs[0]
p.text = "Thank You!"
p.font.size = Pt(56)
p.font.bold = True
p.font.color.rgb = TEXT_WHITE
p.alignment = PP_ALIGN.CENTER

cb = slide.shapes.add_textbox(Inches(0), Inches(3.5), Inches(13.333), Inches(2.2))
tf = cb.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "Damilola Olaiya"
p.font.size = Pt(18)
p.font.color.rgb = ACCENT3
p.alignment = PP_ALIGN.CENTER

p = tf.add_paragraph()
p.text = "damilolaolaiya@cmail.carleton.ca"
p.font.size = Pt(13)
p.font.color.rgb = TEXT_LIGHT
p.alignment = PP_ALIGN.CENTER

p = tf.add_paragraph()
p.text = ""

p = tf.add_paragraph()
p.text = "Mithun Mani"
p.font.size = Pt(18)
p.font.color.rgb = ACCENT3
p.alignment = PP_ALIGN.CENTER

p = tf.add_paragraph()
p.text = "mithunmani@cmail.carleton.ca"
p.font.size = Pt(13)
p.font.color.rgb = TEXT_LIGHT
p.alignment = PP_ALIGN.CENTER

rb = slide.shapes.add_textbox(Inches(0), Inches(6.2), Inches(13.333), Inches(0.5))
p = rb.text_frame.paragraphs[0]
p.text = "github.com/whiteblaze143/DATA_5000"
p.font.size = Pt(16)
p.font.color.rgb = ACCENT
p.alignment = PP_ALIGN.CENTER

# Save
out = "/home/mithunmanivannan/DATA_5000/DATA_5000/ECG_Reconstruction_Final.pptx"
prs.save(out)
print(f"✓ Presentation saved: {out}")
print(f"  Total slides: {len(prs.slides)}")
