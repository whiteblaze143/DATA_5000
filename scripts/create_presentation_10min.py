#!/usr/bin/env python3
"""
10-Minute Presentation for ECG Reconstruction Project
DATA 5000 - Carleton University
Authors: Damilola Olaiya & Mithun Mani

TRIMMED VERSION: 21 slides for 10 min presentation + 5 min Q&A
~30 seconds per slide average

Slide Budget:
- Title: 1 slide (30s)
- Motivation: 2 slides (1 min)
- Methods: 3 slides (1.5 min)
- Data: 1 slide (30s)
- Results: 5 slides (2.5 min)
- Analysis: 3 slides (1.5 min)
- Conclusion: 3 slides (1.5 min)
- Q&A + Thanks: 2 slides (1 min)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os
import json

# Presentation setup
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
PRIMARY = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0xE9, 0x4D, 0x60)
ACCENT2 = RGBColor(0x0F, 0x4C, 0x75)
ACCENT3 = RGBColor(0x3E, 0xBA, 0xC2)
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_LIGHT = RGBColor(0xE0, 0xE0, 0xE0)
SUCCESS = RGBColor(0x4E, 0xC9, 0xB0)
WARNING = RGBColor(0xFF, 0xCC, 0x00)

FIGURES = "/home/mithunmanivannan/DATA_5000/DATA_5000/docs/figures"
MODELS = "/home/mithunmanivannan/DATA_5000/DATA_5000/models"


# =============================================================================
# PLACEHOLDER: Load actual results when available
# =============================================================================
def load_results():
    """Load latest training results. Returns placeholder if not available."""
    results_path = f"{MODELS}/final_exp_baseline/test_results.json"
    if os.path.exists(results_path):
        with open(results_path) as f:
            return json.load(f)
    
    # Placeholder values - UPDATE AFTER TRAINING
    return {
        'test_correlation_overall': 0.927,  # UPDATE
        'test_correlation_per_lead': [1.0, 1.0, 1.0, 1.0, 1.0, 1.0, 
                                       0.78, 0.72, 0.80, 1.0, 0.87, 0.78],  # UPDATE
        'test_mae_overall': 0.012,  # UPDATE
        'test_snr_overall': 65.0,  # UPDATE
        'dl_leads_corr': 0.82,  # UPDATE
    }

RESULTS = load_results()


# =============================================================================
# HELPER FUNCTIONS
# =============================================================================
def bg(slide, color=PRIMARY):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    spTree = slide.shapes._spTree
    sp = s._element
    spTree.remove(sp)
    spTree.insert(2, sp)

def bar(slide, y=0, h=0.06, color=ACCENT):
    b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(y), Inches(13.333), Inches(h))
    b.fill.solid()
    b.fill.fore_color.rgb = color
    b.line.fill.background()

def title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    bar(slide, 0, 0.12, ACCENT)
    
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(2.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.alignment = PP_ALIGN.CENTER
    
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(4.1), Inches(4.3), Inches(0.02)).fill.solid()
    
    sb = slide.shapes.add_textbox(Inches(0.8), Inches(4.3), Inches(11.7), Inches(0.8))
    p = sb.text_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(20)
    p.font.color.rgb = ACCENT3
    p.alignment = PP_ALIGN.CENTER
    
    ab = slide.shapes.add_textbox(Inches(0.8), Inches(5.4), Inches(11.7), Inches(0.5))
    p = ab.text_frame.paragraphs[0]
    p.text = "Damilola Olaiya  •  Mithun Mani"
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_LIGHT
    p.alignment = PP_ALIGN.CENTER
    
    cb = slide.shapes.add_textbox(Inches(0.8), Inches(6.1), Inches(11.7), Inches(0.4))
    p = cb.text_frame.paragraphs[0]
    p.text = "DATA 5000  |  Carleton University  |  December 2025"
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_LIGHT
    p.alignment = PP_ALIGN.CENTER
    
    bar(slide, 7.38, 0.12, ACCENT)

def section_slide(prs, num, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, ACCENT)
    
    nb = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(11.7), Inches(1.5))
    p = nb.text_frame.paragraphs[0]
    p.text = f"0{num}" if num < 10 else str(num)
    p.font.size = Pt(120)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.alignment = PP_ALIGN.CENTER
    
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(11.7), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = title.upper()
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.alignment = PP_ALIGN.CENTER

def content_slide(prs, title, bullets, fig=None, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    bar(slide, 0, 0.06, ACCENT)
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.65))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    
    cw = 6 if fig and os.path.exists(fig) else 12
    cb = slide.shapes.add_textbox(Inches(0.5), Inches(0.95), Inches(cw), Inches(6))
    tf = cb.text_frame
    tf.word_wrap = True
    
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        txt = b.strip().lstrip("•-▸ ")
        indent = 1 if b.startswith("  ") else 0
        
        if txt.startswith("✓"): p.font.color.rgb = SUCCESS
        elif txt.startswith("△") or txt.startswith("⚠"): p.font.color.rgb = WARNING
        elif txt.startswith("✗"): p.font.color.rgb = ACCENT
        else: p.font.color.rgb = TEXT_LIGHT
        
        p.text = ("▸ " + txt) if indent == 0 and txt else ("   " + txt)
        p.level = indent
        p.font.size = Pt(16 if indent == 0 else 14)
        p.space_after = Pt(5)
    
    if fig and os.path.exists(fig):
        slide.shapes.add_picture(fig, Inches(6.7), Inches(1.1), width=Inches(6.2))
    
    if notes:
        slide.notes_slide.notes_text_frame.text = notes

def figure_slide(prs, title, fig, caption=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    bar(slide, 0, 0.06, ACCENT)
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    
    if os.path.exists(fig):
        slide.shapes.add_picture(fig, Inches(0.7), Inches(0.9), width=Inches(11.9))
    else:
        # Placeholder
        placeholder = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(9), Inches(2))
        p = placeholder.text_frame.paragraphs[0]
        p.text = f"[Figure: {os.path.basename(fig)}]\nGenerate after training completes"
        p.font.size = Pt(20)
        p.font.color.rgb = WARNING
        p.alignment = PP_ALIGN.CENTER
    
    if caption:
        cb = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12), Inches(0.5))
        p = cb.text_frame.paragraphs[0]
        p.text = caption
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_LIGHT
        p.alignment = PP_ALIGN.CENTER

def table_slide(prs, title, headers, rows, highlight=-1):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    bar(slide, 0, 0.06, ACCENT)
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    
    nr, nc = len(rows) + 1, len(headers)
    tw = min(12, nc * 2.1)
    cw = tw / nc
    
    left = Inches((13.333 - tw) / 2)
    table = slide.shapes.add_table(nr, nc, left, Inches(1.1), Inches(tw), Inches(0.4 * nr)).table
    
    for i in range(nc):
        table.columns[i].width = Inches(cw)
    
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_WHITE
        p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            cell.fill.solid()
            if ri == highlight:
                cell.fill.fore_color.rgb = RGBColor(0x1A, 0x4A, 0x3A)
            else:
                cell.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x3A) if ri % 2 == 0 else RGBColor(0x1C, 0x1C, 0x32)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.font.color.rgb = TEXT_LIGHT
            p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

def metrics_slide(prs, title, metrics):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    bar(slide, 0, 0.06, ACCENT)
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    
    n = len(metrics)
    cw = 3.0
    sp = (13.333 - n * cw) / (n + 1)
    
    for i, (label, val, color) in enumerate(metrics):
        x = sp + i * (cw + sp)
        
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.6), Inches(cw), Inches(4))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0x20, 0x20, 0x38)
        card.line.color.rgb = color
        card.line.width = Pt(2)
        
        vb = slide.shapes.add_textbox(Inches(x), Inches(2.4), Inches(cw), Inches(1.2))
        p = vb.text_frame.paragraphs[0]
        p.text = val
        p.font.size = Pt(42)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER
        
        lb = slide.shapes.add_textbox(Inches(x), Inches(3.8), Inches(cw), Inches(1.4))
        tf = lb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_LIGHT
        p.alignment = PP_ALIGN.CENTER


# =============================================================================
# BUILD 21-SLIDE PRESENTATION
# =============================================================================

# SLIDE 1: Title (30s)
title_slide(prs, 
    "Efficient 12-Lead ECG Reconstruction\nfrom Reduced Lead Sets",
    "A Computationally Lightweight Physics-Informed Approach")

# SLIDE 2: The Problem (45s)
content_slide(prs, "The Clinical Need", [
    "Cardiovascular Disease: #1 cause of death globally",
    "  17.9 million deaths annually",
    "",
    "12-Lead ECG: Gold standard diagnosis",
    "  But requires 10 electrodes + trained personnel",
    "",
    "The Gap:",
    "  Consumer wearables: only 1-3 leads",
    "  Full diagnosis needs all 12 leads",
    "",
    "Our Goal: Reconstruct 12 leads from 3"
], fig=f"{FIGURES}/problem_visualization.png")

# SLIDE 3: Our Key Insight (45s)
content_slide(prs, "Key Insight: Physics Reduces the Problem", [
    "The 12 leads are NOT independent!",
    "",
    "✓ 4 Leads are EXACT (Physics):",
    "  III = II - I     (Einthoven 1903)",
    "  aVR, aVL, aVF    (Goldberger)",
    "",
    "✓ Only 5 Leads Need Deep Learning:",
    "  V1, V2, V3, V5, V6",
    "",
    "Result:",
    "  Zero parameters for 4 leads",
    "  Perfect reconstruction guaranteed",
    "  6x fewer params than SOTA"
], fig=f"{FIGURES}/approach_diagram.png")

# SLIDE 4: Architecture (45s)
content_slide(prs, "Hybrid Architecture", [
    "Input: 3 leads (I, II, V4)",
    "",
    "Physics Module:",
    "  Einthoven + Goldberger laws",
    "  → 4 leads, r = 1.0, zero cost",
    "",
    "Deep Learning Module:",
    "  1D U-Net (encoder-decoder)",
    "  17.1M parameters",
    "  → 5 chest leads",
    "",
    "Total: Complete 12-lead ECG"
], fig=f"{FIGURES}/architecture_diagram_clean.png")

# SLIDE 5: Experimental Rigor (30s)
content_slide(prs, "Rigorous Methodology", [
    "Frozen Hyperparameters (LR-validated):",
    "  lr = 3×10⁻⁴, batch = 128, epochs = 150",
    "  AdamW + ReduceLROnPlateau",
    "",
    "Patient-Wise Data Splits:",
    "  No data leakage between train/test",
    "  Stricter than most prior work",
    "",
    "Statistical Framework:",
    "  Paired t-test, Cohen's d effect size",
    "  Bootstrap 95% CI"
])

# SLIDE 6: Dataset (30s)
content_slide(prs, "Dataset: PTB-XL", [
    "21,837 clinical 12-lead ECGs",
    "18,885 unique patients",
    "500 Hz, 10 seconds",
    "",
    "Train: 14,363 | Val: 1,914 | Test: 1,932",
    "",
    "Preprocessing:",
    "  Z-score normalization",
    "  Patient-wise stratified splits"
], fig=f"{FIGURES}/dataset_characteristics.png")

# SLIDE 7: Training Convergence (30s)
figure_slide(prs, "Training Convergence",
    f"{FIGURES}/training_curves_final.png",
    "Smooth convergence with ReduceLROnPlateau scheduler")

# SLIDE 8: Key Results - Metrics (45s)
# Use placeholder values - will update after training
dl_corr = f"r = {RESULTS.get('dl_leads_corr', 0.82):.2f}"
overall_corr = f"r = {RESULTS.get('test_correlation_overall', 0.927):.2f}"

metrics_slide(prs, "Overall Performance", [
    ("Physics Leads\n(III, aVR, aVL, aVF)", "r = 1.0", SUCCESS),
    ("DL Leads\n(V1-V3, V5-V6)", dl_corr, ACCENT3),
    ("Overall\n12-Lead", overall_corr, SUCCESS),
])

# SLIDE 9: Per-Lead Results (45s)
figure_slide(prs, "Per-Lead Performance",
    f"{FIGURES}/per_lead_barplot.png",
    "Physics leads perfect, DL leads V5 best (closest to V4), V2 hardest")

# SLIDE 10: Sample Reconstructions (30s)
figure_slide(prs, "Sample ECG Reconstructions",
    f"{FIGURES}/reconstruction_samples.png",
    "Visual comparison: Ground truth (blue) vs Reconstructed (red)")

# SLIDE 11: Why Shared Decoder Wins (45s)
content_slide(prs, "Ablation: Shared vs Lead-Specific Decoder", [
    "Surprising Result:",
    "  Shared (17.1M) BEATS Lead-Specific (40.8M)",
    "",
    "Statistical Evidence:",
    "  Cohen's d = 0.92 (large effect)",
    "  Bootstrap 95% CI: [0.006, 0.072]",
    "  → Excludes zero = significant",
    "",
    "Interpretation:",
    "  Limited input (3 leads) = limited information",
    "  More params → overfitting, not capacity",
    "  Parameter sharing = beneficial regularization"
], fig=f"{FIGURES}/model_comparison_detailed.png")

# SLIDE 12: Information Bottleneck (45s)
content_slide(prs, "Why V2 is Hardest: Information Bottleneck", [
    "Ground truth correlation V4 ↔ V2 = 0.36",
    "  Fundamental limit, not model failure",
    "",
    "Best performance on V5 (r ≈ 0.87):",
    "  V4 ↔ V5 correlation = 0.79",
    "",
    "Worst performance on V2 (r ≈ 0.72):",
    "  V2 anatomically distant from V4",
    "",
    "Implication:",
    "  Input lead selection matters more",
    "  than model architecture"
], fig=f"{FIGURES}/lead_correlation_heatmap.png")

# SLIDE 13: Comparison with SOTA (45s)
table_slide(prs, "Comparison with State-of-the-Art",
    ["Method", "Year", "DL Leads r", "Params", "Split"],
    [["Linear (Frank)", "1970s", "0.70-0.75", "~0", "N/A"],
     ["CNN (Sohn)", "2020", "~0.85", "~30M", "Record"],
     ["LSTM (Lee)", "2021", "~0.88", "~60M", "Record"],
     ["Transformer", "2023", "~0.90", "100M+", "Record"],
     ["Ours", "2025", "~0.82", "17.1M", "Patient"]],  # UPDATE
    highlight=4)

# SLIDE 14: Honest Assessment (30s)
content_slide(prs, "Honest Assessment", [
    "Our Advantages:",
    "  ✓ 6x fewer parameters (17M vs 100M)",
    "  ✓ Physics leads perfect (r = 1.0)",
    "  ✓ Stricter evaluation (patient-split)",
    "",
    "Our Limitations:",
    "  △ DL leads slightly below SOTA",
    "  △ Single dataset (PTB-XL only)",
    "  △ Input config (I, II, V4) not optimized",
    "",
    "Key Insight:",
    "  Gap may be due to stricter evaluation"
])

# SLIDE 15: Key Contributions (45s)
content_slide(prs, "Key Contributions", [
    "1. Efficient Hybrid Architecture",
    "   Physics + lightweight DL = 6x smaller",
    "",
    "2. Rigorous Methodology",
    "   Patient-wise splits, frozen HPs, stats tests",
    "",
    "3. Information Bottleneck Analysis",
    "   Identified fundamental reconstruction limits",
    "",
    "4. Open Source",
    "   github.com/whiteblaze143/DATA_5000"
])

# SLIDE 16: Limitations (30s)
content_slide(prs, "Limitations & Future Work", [
    "Current Limitations:",
    "  Single dataset (PTB-XL)",
    "  Input config not optimized",
    "  No downstream task validation",
    "",
    "Future Directions:",
    "  Test (I, II, V2) or (I, II, V2, V4)",
    "  External validation (Chapman, MIMIC)",
    "  Multi-label classification on reconstructed ECGs",
    "  Uncertainty quantification"
])

# SLIDE 17: Summary Metrics (30s)
metrics_slide(prs, "Summary", [
    ("Parameters", "17.1M\n(6x smaller)", SUCCESS),
    ("Physics Leads", "r = 1.0\n(guaranteed)", SUCCESS),
    ("DL Leads", "r ≈ 0.82\n(patient-split)", ACCENT3),
])

# SLIDE 18: Code & Resources (15s)
content_slide(prs, "Code & Resources", [
    "GitHub Repository:",
    "  github.com/whiteblaze143/DATA_5000",
    "",
    "Key Files:",
    "  run_training.py — Main entry",
    "  src/physics.py — Einthoven/Goldberger",
    "  src/models/unet_1d.py — 1D U-Net",
    "",
    "Trained Models & Report:",
    "  models/final_exp_baseline/",
    "  docs/PROJECT_REPORT.pdf"
])

# SLIDE 19: Acknowledgments (15s)
content_slide(prs, "Acknowledgments", [
    "Course: DATA 5000 - Data Science Capstone",
    "",
    "Carleton University",
    "",
    "PhysioNet for PTB-XL dataset",
    "",
    "Google Cloud for compute resources"
])

# SLIDE 20: Questions
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)

qb = slide.shapes.add_textbox(Inches(0), Inches(1.3), Inches(13.333), Inches(3))
p = qb.text_frame.paragraphs[0]
p.text = "?"
p.font.size = Pt(200)
p.font.bold = True
p.font.color.rgb = ACCENT
p.alignment = PP_ALIGN.CENTER

tb = slide.shapes.add_textbox(Inches(0), Inches(4.5), Inches(13.333), Inches(1))
p = tb.text_frame.paragraphs[0]
p.text = "Questions?"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = TEXT_WHITE
p.alignment = PP_ALIGN.CENTER

# SLIDE 21: Thank You
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
bar(slide, 0, 0.1, ACCENT)
bar(slide, 7.4, 0.1, ACCENT)

tyb = slide.shapes.add_textbox(Inches(0), Inches(1.8), Inches(13.333), Inches(1.2))
p = tyb.text_frame.paragraphs[0]
p.text = "Thank You!"
p.font.size = Pt(56)
p.font.bold = True
p.font.color.rgb = TEXT_WHITE
p.alignment = PP_ALIGN.CENTER

cb = slide.shapes.add_textbox(Inches(0), Inches(3.5), Inches(13.333), Inches(2.2))
tf = cb.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "Damilola Olaiya"
p.font.size = Pt(18)
p.font.color.rgb = ACCENT3
p.alignment = PP_ALIGN.CENTER

p = tf.add_paragraph()
p.text = "damilolaolaiya@cmail.carleton.ca"
p.font.size = Pt(13)
p.font.color.rgb = TEXT_LIGHT
p.alignment = PP_ALIGN.CENTER

p = tf.add_paragraph()
p.text = ""

p = tf.add_paragraph()
p.text = "Mithun Mani"
p.font.size = Pt(18)
p.font.color.rgb = ACCENT3
p.alignment = PP_ALIGN.CENTER

p = tf.add_paragraph()
p.text = "mithunmani@cmail.carleton.ca"
p.font.size = Pt(13)
p.font.color.rgb = TEXT_LIGHT
p.alignment = PP_ALIGN.CENTER

rb = slide.shapes.add_textbox(Inches(0), Inches(6.2), Inches(13.333), Inches(0.5))
p = rb.text_frame.paragraphs[0]
p.text = "github.com/whiteblaze143/DATA_5000"
p.font.size = Pt(16)
p.font.color.rgb = ACCENT
p.alignment = PP_ALIGN.CENTER


# =============================================================================
# SAVE
# =============================================================================
out = "/home/mithunmanivannan/DATA_5000/DATA_5000/ECG_Reconstruction_10min.pptx"
prs.save(out)
print(f"✓ Presentation saved: {out}")
print(f"  Total slides: {len(prs.slides)}")
print(f"\n  Slide breakdown:")
print(f"    1. Title")
print(f"    2-3. Motivation (Problem + Key Insight)")
print(f"    4-5. Methods (Architecture + Rigor)")
print(f"    6. Dataset")
print(f"    7-10. Results (Training, Metrics, Per-Lead, Samples)")
print(f"    11-14. Analysis (Ablation, Bottleneck, SOTA, Honest)")
print(f"    15-17. Conclusion (Contributions, Limitations, Summary)")
print(f"    18-19. Resources + Acknowledgments")
print(f"    20-21. Q&A + Thank You")
print(f"\n  Time budget: ~30s per slide = 10 min + 5 min Q&A")
print(f"\n  ⚠ UPDATE RESULTS after training completes!")
print(f"    Run: python scripts/generate_final_figures.py --model_dir models/final_exp_baseline")
print(f"    Then re-run this script to update metrics.")
